from diffusers import QwenImageLayeredPipeline
import torch
from PIL import Image
from pptx import Presentation
import os
import gradio as gr
import uuid
import numpy as np
import random
import tempfile
import zipfile 

MAX_SEED = np.iinfo(np.int32).max

pipeline = QwenImageLayeredPipeline.from_pretrained("Qwen/Qwen-Image-Layered")
pipeline = pipeline.to("cuda", torch.bfloat16)
pipeline.set_progress_bar_config(disable=None)

def imagelist_to_pptx(img_files):
    with Image.open(img_files[0]) as img:
        img_width_px, img_height_px = img.size

    def px_to_emu(px, dpi=96):
        inch = px / dpi
        emu = inch * 914400
        return int(emu)

    prs = Presentation()
    prs.slide_width = px_to_emu(img_width_px)
    prs.slide_height = px_to_emu(img_height_px)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = top = 0
    for img_path in img_files:
        slide.shapes.add_picture(img_path, left, top, width=px_to_emu(img_width_px), height=px_to_emu(img_height_px))

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        return tmp.name

def export_gallery(images):
    # images: list of image file paths
    images = [e[0] for e in images]
    pptx_path = imagelist_to_pptx(images)
    return pptx_path

def infer(input_image,
          seed=777,
          randomize_seed=False,
          prompt=None,
          neg_prompt=" ",
          true_guidance_scale=4.0,
          num_inference_steps=50,
          layer=4,
          cfg_norm=True,
          use_en_prompt=True):
    
    if randomize_seed:
        seed = random.randint(0, MAX_SEED)
        
    if isinstance(input_image, list):
        input_image = input_image[0]
        
    if isinstance(input_image, str):
        pil_image = Image.open(input_image).convert("RGB").convert("RGBA")
    elif isinstance(input_image, Image.Image):
        pil_image = input_image.convert("RGB").convert("RGBA")
    elif isinstance(input_image, np.ndarray):
        pil_image = Image.fromarray(input_image).convert("RGB").convert("RGBA")
    else:
        raise ValueError("Unsupported input_image type: %s" % type(input_image))
    
    inputs = {
        "image": pil_image,
        "generator": torch.Generator(device='cuda').manual_seed(seed),
        "true_cfg_scale": true_guidance_scale,
        "prompt": prompt,
        "negative_prompt": neg_prompt,
        "num_inference_steps": num_inference_steps,
        "num_images_per_prompt": 1,
        "layers": layer,
        "resolution": 640,      # Using different bucket (640, 1024) to determine the resolution. For this version, 640 is recommended
        "cfg_normalize": cfg_norm,  # Whether enable cfg normalization.
        "use_en_prompt": use_en_prompt, 
    }
    print(inputs)
    with torch.inference_mode():
        output = pipeline(**inputs)
        output_images = output.images[0]
    
    output = []
    temp_files = []
    for i, image in enumerate(output_images):
        output.append(image)
        # Save to temp file for export
        tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        image.save(tmp.name)
        temp_files.append(tmp.name)
    
    # Generate PPTX
    pptx_path = imagelist_to_pptx(temp_files)
    
    # Generate ZIP
    with tempfile.NamedTemporaryFile(suffix=".zip", delete=False) as tmp:
        with zipfile.ZipFile(tmp.name, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for i, img_path in enumerate(temp_files):
                zipf.write(img_path, f"layer_{i+1}.png")
        zip_path = tmp.name
    
    return output, pptx_path, zip_path

examples = [
            "assets/test_images/1.png",
            "assets/test_images/2.png",
            "assets/test_images/3.png",
            "assets/test_images/4.png",
            "assets/test_images/5.png",
            "assets/test_images/6.png",
            "assets/test_images/7.png",
            "assets/test_images/8.png",
            "assets/test_images/9.png",
            "assets/test_images/10.png",
            "assets/test_images/11.png",
            "assets/test_images/12.png",
            "assets/test_images/13.png",
            ]


with gr.Blocks() as demo:
    with gr.Column(elem_id="col-container"):
        gr.HTML('''<p align="center"><img src="https://qianwen-res.oss-cn-beijing.aliyuncs.com/Qwen-Image/layered/qwen-image-layered-logo.png" width="400"/><p>''')
        gr.Markdown("""
                    The text prompt is intended to describe the overall content of the input image—including elements that may be partially occluded (e.g., you may specify the text hidden behind a foreground object). It is not designed to control the semantic content of individual layers explicitly.
                    """)
        with gr.Row():
            with gr.Column(scale=1):
                input_image = gr.Image(label="Input Image", image_mode="RGBA")
                
                
                with gr.Accordion("Advanced Settings", open=False):
                    prompt = gr.Textbox(
                        label="Prompt (Optional)",
                        placeholder="Please enter the prompt to describe the image. It is not designed to control the semantic content of individual layers explicitly. (Optional)",
                        value="",
                        lines=3,
                    )
                    neg_prompt = gr.Textbox(
                        label="Negative Prompt (Optional)",
                        placeholder="Please enter the negative prompt",
                        value=" ",
                        lines=3,
                    )
                    
                    seed = gr.Slider(
                        label="Seed",
                        minimum=0,
                        maximum=MAX_SEED,
                        step=1,
                        value=0,
                    )
                    randomize_seed = gr.Checkbox(label="Randomize seed", value=True)
                    
                    true_guidance_scale = gr.Slider(
                        label="True guidance scale",
                        minimum=1.0,
                        maximum=10.0,
                        step=0.1,
                        value=4.0
                    )

                    num_inference_steps = gr.Slider(
                        label="Number of inference steps",
                        minimum=1,
                        maximum=50,
                        step=1,
                        value=50,
                    )

                    layer = gr.Slider(
                        label="Layers",
                        minimum=2,
                        maximum=10,
                        step=1,
                        value=4,
                    )

                    cfg_norm = gr.Checkbox(label="Whether enable CFG normalization", value=True)
                    use_en_prompt = gr.Checkbox(label="Automatic caption language if no prompt provided, True for EN, False for ZH", value=True)
                
                run_button = gr.Button("Decompose!", variant="primary")

            with gr.Column(scale=2):
                gallery = gr.Gallery(label="Layers", columns=4, rows=1, format="png")
                with gr.Row():
                    export_file = gr.File(label="Download PPTX")
                    export_zip_file = gr.File(label="Download ZIP")

    gr.Examples(examples=examples,
                inputs=[input_image], 
                outputs=[gallery, export_file, export_zip_file],
                fn=infer, 
                examples_per_page=14,
                cache_examples=False,
                run_on_click=True
    )

    run_button.click(
        fn=infer,
        inputs=[
            input_image,
            seed,
            randomize_seed,
            prompt,
            neg_prompt,
            true_guidance_scale,
            num_inference_steps,
            layer,
            cfg_norm,
            use_en_prompt,
        ], 
        outputs=[gallery, export_file, export_zip_file],
    )

demo.launch(
    server_name="0.0.0.0",
    server_port=7869,
)
